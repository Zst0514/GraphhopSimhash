#!/usr/bin/env python3
"""Draw the GRACE introduction overview PPT."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "HPCA_2027_GFMAcc" / "Figure" / "intro_overview_grace_v2.pptx"


COL = {
    "ink": RGBColor(17, 24, 39),
    "muted": RGBColor(71, 84, 103),
    "blue": RGBColor(219, 234, 254),
    "blue_edge": RGBColor(37, 99, 235),
    "green": RGBColor(220, 252, 231),
    "green_edge": RGBColor(22, 163, 74),
    "orange": RGBColor(255, 237, 213),
    "orange_edge": RGBColor(234, 88, 12),
    "purple": RGBColor(237, 233, 254),
    "purple_edge": RGBColor(124, 58, 237),
    "gray": RGBColor(241, 245, 249),
    "gray_edge": RGBColor(148, 163, 184),
    "red": RGBColor(254, 226, 226),
    "red_edge": RGBColor(220, 38, 38),
    "yellow": RGBColor(254, 243, 199),
    "yellow_edge": RGBColor(217, 119, 6),
    "white": RGBColor(255, 255, 255),
}


def set_text(shape, text, size=14, bold=False, color="ink", align=PP_ALIGN.CENTER):
    shape.text = text
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = "Arial"
    run.font.color.rgb = COL[color]
    tf.margin_left = Inches(0.04)
    tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)


def box(slide, x, y, w, h, text="", fill="gray", edge="gray_edge",
        size=13, bold=True, radius=True, align=PP_ALIGN.CENTER):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COL[fill]
    shp.line.color.rgb = COL[edge]
    shp.line.width = Pt(1.5)
    if text:
        set_text(shp, text, size=size, bold=bold, align=align)
    return shp


def text(slide, x, y, w, h, s, size=13, bold=False, color="ink",
         align=PP_ALIGN.LEFT):
    shp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shp, s, size=size, bold=bold, color=color, align=align)
    return shp


def line(slide, x1, y1, x2, y2, color="ink", width=1.6, arrow=True):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = COL[color]
    conn.line.width = Pt(width)
    if arrow:
        conn.line.end_arrowhead = True
    return conn


def panel(slide, x, y, w, h, title):
    box(slide, x, y, w, h, fill="white", edge="ink", size=10, bold=False, radius=True)
    text(slide, x + 0.12, y + 0.08, w - 0.24, 0.30, title, size=15, bold=True,
         align=PP_ALIGN.CENTER)


def graph_icon(slide, x, y, scale=1.0):
    pts = [(x, y + 0.35), (x + 0.45, y), (x + 0.9, y + 0.32),
           (x + 0.35, y + 0.78), (x + 0.82, y + 0.88)]
    edges = [(0, 1), (1, 2), (0, 3), (3, 2), (3, 4), (2, 4)]
    for a, b in edges:
        line(slide, pts[a][0], pts[a][1], pts[b][0], pts[b][1], color="blue_edge",
             width=1.2, arrow=False)
    for px, py in pts:
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(px - 0.055 * scale),
                                     Inches(py - 0.055 * scale),
                                     Inches(0.11 * scale), Inches(0.11 * scale))
        shp.fill.solid()
        shp.fill.fore_color.rgb = COL["white"]
        shp.line.color.rgb = COL["blue_edge"]
        shp.line.width = Pt(1.4)


def mini_array(slide, x, y, w=0.70, h=0.54, rows=2, cols=3, fill="white"):
    box(slide, x, y, w, h, fill=fill, edge="ink", size=6, bold=False, radius=False)
    cw = w / cols
    ch = h / rows
    for r in range(rows):
        for c in range(cols):
            shp = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + c * cw + 0.03),
                Inches(y + r * ch + 0.03),
                Inches(cw - 0.06),
                Inches(ch - 0.06),
            )
            shp.fill.solid()
            shp.fill.fore_color.rgb = COL["white"]
            shp.line.color.rgb = COL["ink"]
            shp.line.width = Pt(0.8)


def draw():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    text(slide, 0.25, 0.08, 12.85, 0.42,
         "GRACE: Graph-Aware Reuse and Progressive BFP Execution for TAG Encoders",
         size=20, bold=True, align=PP_ALIGN.CENTER)

    # Panels.
    panel(slide, 0.18, 0.62, 3.95, 6.40, "(a) TAG LLM Encoder Bottleneck")
    panel(slide, 4.30, 0.62, 4.55, 6.40, "(b) Workload Observations")
    panel(slide, 9.02, 0.62, 4.13, 6.40, "(c) GRACE Co-Design")

    # Main arrows between panels.
    line(slide, 4.12, 3.55, 4.30, 3.55, color="purple_edge", width=3.0)
    line(slide, 8.85, 3.55, 9.02, 3.55, color="purple_edge", width=3.0)

    # Panel A: pipeline.
    graph_icon(slide, 0.55, 1.30, scale=1.3)
    text(slide, 0.42, 2.35, 1.25, 0.38, "TAG nodes", size=12, bold=True,
         align=PP_ALIGN.CENTER)
    box(slide, 1.75, 1.35, 1.92, 0.82, "LLM / TF\nencoder", fill="blue",
        edge="blue_edge", size=13)
    box(slide, 1.98, 2.70, 1.48, 0.66, "GNN\nclassifier", fill="green",
        edge="green_edge", size=12)
    line(slide, 1.47, 1.78, 1.75, 1.78, color="ink", width=1.8)
    line(slide, 2.70, 2.17, 2.70, 2.70, color="ink", width=1.8)
    line(slide, 3.46, 3.03, 3.80, 3.03, color="ink", width=1.8)
    text(slide, 0.55, 4.05, 3.20, 0.40, "End-to-end cost is dominated by\nnode-wise LLM encoding",
         size=13, bold=True, align=PP_ALIGN.CENTER)
    box(slide, 0.72, 4.82, 2.82, 0.44, "LLM encoder", fill="blue",
        edge="blue_edge", size=11)
    box(slide, 0.72, 5.34, 0.52, 0.36, "GNN", fill="green", edge="green_edge",
        size=9)
    text(slide, 1.38, 5.30, 2.20, 0.48, "small backend cost,\nlarge graph effect",
         size=10.5, color="muted")
    text(slide, 0.55, 6.18, 3.05, 0.48,
         "Ordinary TF accelerators treat nodes as\nindependent dense text sequences.",
         size=10.5, color="muted", align=PP_ALIGN.CENTER)

    # Panel B: observations.
    box(slide, 4.55, 1.18, 4.05, 1.42, fill="blue", edge="blue_edge")
    text(slide, 4.78, 1.30, 3.58, 0.28, "1. Graph locality exposes reuse", size=13,
         bold=True)
    graph_icon(slide, 4.78, 1.78, scale=0.9)
    text(slide, 5.90, 1.70, 2.30, 0.52,
         "Neighbor nodes are closer\nthan random pairs", size=10.5)
    box(slide, 7.55, 2.00, 0.75, 0.28, "~34%", fill="white", edge="blue_edge",
        size=10)
    text(slide, 7.42, 2.28, 1.05, 0.20, "low-drop reuse", size=8.5,
         color="muted", align=PP_ALIGN.CENTER)

    box(slide, 4.55, 2.88, 4.05, 1.44, fill="orange", edge="orange_edge")
    text(slide, 4.78, 3.00, 3.58, 0.28, "2. Encoder precision has a boundary",
         size=13, bold=True)
    box(slide, 4.84, 3.55, 0.70, 0.34, "BFPA3\nfails", fill="red",
        edge="red_edge", size=8.5)
    box(slide, 5.80, 3.48, 0.82, 0.46, "BFPA4\nbase", fill="yellow",
        edge="yellow_edge", size=9.5)
    box(slide, 6.88, 3.42, 0.82, 0.58, "BFPA6\nsafe lift", fill="green",
        edge="green_edge", size=9.5)
    line(slide, 5.54, 3.72, 5.80, 3.72, color="orange_edge", width=1.4)
    line(slide, 6.62, 3.72, 6.88, 3.72, color="orange_edge", width=1.4)
    text(slide, 7.88, 3.38, 0.58, 0.64, "block\nselective", size=9.0,
         color="muted", align=PP_ALIGN.CENTER)

    box(slide, 4.55, 4.58, 4.05, 1.50, fill="purple", edge="purple_edge")
    text(slide, 4.78, 4.70, 3.58, 0.28, "3. Filtering and GEMMs diverge", size=13,
         bold=True)
    box(slide, 4.88, 5.28, 1.18, 0.44, "reuse\nfilter", fill="gray",
        edge="gray_edge", size=9.5)
    box(slide, 6.80, 5.18, 1.28, 0.62, "dense TF\nGEMMs", fill="blue",
        edge="blue_edge", size=9.5)
    line(slide, 6.06, 5.50, 6.80, 5.50, color="purple_edge", width=1.5)
    text(slide, 4.83, 5.82, 3.40, 0.24, "low-intensity state access vs. high-intensity tensor compute",
         size=8.7, color="muted", align=PP_ALIGN.CENTER)

    # Panel C: GRACE framework.
    box(slide, 9.32, 1.12, 1.72, 4.90, "Near-DIMM\nreuse frontend", fill="green",
        edge="green_edge", size=12)
    box(slide, 9.52, 1.80, 1.32, 0.50, "graph-context\nlookup", fill="white",
        edge="green_edge", size=8.8)
    box(slide, 9.52, 2.55, 1.32, 0.50, "graph-aware\nreuse filter", fill="white",
        edge="green_edge", size=8.8)
    box(slide, 9.52, 3.30, 1.32, 0.50, "prior encoder\nresults", fill="white",
        edge="green_edge", size=8.8)
    box(slide, 9.52, 4.05, 1.32, 0.50, "computed-node\nselection", fill="white",
        edge="green_edge", size=8.8)

    box(slide, 11.34, 1.12, 1.54, 4.90, "BFP NPU\nbackend", fill="blue",
        edge="blue_edge", size=12)
    box(slide, 11.55, 1.88, 1.10, 0.42, "BFPA4 base", fill="yellow",
        edge="yellow_edge", size=9.3)
    box(slide, 11.55, 2.52, 1.10, 0.42, "BFPA6 lift", fill="orange",
        edge="orange_edge", size=9.3)
    mini_array(slide, 11.64, 3.22, w=0.92, h=0.62, rows=2, cols=3)
    box(slide, 11.55, 4.35, 1.10, 0.44, "TF GEMMs", fill="white",
        edge="blue_edge", size=9.3)

    line(slide, 11.04, 3.28, 11.34, 3.28, color="ink", width=2.0)
    text(slide, 10.85, 3.45, 0.68, 0.42, "nodes +\ngraph tags", size=8.0,
         color="muted", align=PP_ALIGN.CENTER)
    line(slide, 12.88, 3.28, 13.10, 3.28, color="ink", width=2.0)
    text(slide, 12.82, 3.52, 0.42, 0.28, "emb.", size=8.5, color="muted",
         align=PP_ALIGN.CENTER)

    # Contribution tags.
    box(slide, 9.28, 6.20, 3.55, 0.42, "1 reuse before LLM   2 progressive BFP   3 heterogeneous execution",
        fill="purple", edge="purple_edge", size=9.6)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)


if __name__ == "__main__":
    draw()
