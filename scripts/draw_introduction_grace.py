#!/usr/bin/env python3
"""Regenerate Figure/Introduction.pptx for the current GRACE introduction."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "HPCA_2027_GFMAcc" / "Figure" / "Introduction.pptx"


COL = {
    "ink": RGBColor(20, 29, 43),
    "muted": RGBColor(79, 91, 110),
    "blue": RGBColor(219, 234, 254),
    "blue_edge": RGBColor(29, 78, 216),
    "green": RGBColor(220, 252, 231),
    "green_edge": RGBColor(22, 130, 65),
    "orange": RGBColor(255, 237, 213),
    "orange_edge": RGBColor(194, 65, 12),
    "purple": RGBColor(237, 233, 254),
    "purple_edge": RGBColor(109, 40, 217),
    "red": RGBColor(254, 226, 226),
    "red_edge": RGBColor(190, 18, 60),
    "gray": RGBColor(241, 245, 249),
    "gray_edge": RGBColor(100, 116, 139),
    "yellow": RGBColor(254, 243, 199),
    "yellow_edge": RGBColor(180, 83, 9),
    "white": RGBColor(255, 255, 255),
}


def fmt(shape, text, size=6.5, bold=False, color="ink", align=PP_ALIGN.CENTER):
    shape.text = text
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = COL[color]


def box(slide, x, y, w, h, text="", fill="gray", edge="gray_edge",
        size=6.5, bold=True, radius=True, align=PP_ALIGN.CENTER, lw=1.0):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COL[fill]
    shp.line.color.rgb = COL[edge]
    shp.line.width = Pt(lw)
    if text:
        fmt(shp, text, size=size, bold=bold, align=align)
    return shp


def text(slide, x, y, w, h, content, size=6.5, bold=False,
         color="ink", align=PP_ALIGN.LEFT):
    shp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    fmt(shp, content, size=size, bold=bold, color=color, align=align)
    return shp


def line(slide, x1, y1, x2, y2, color="ink", width=1.0, arrow=True):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = COL[color]
    conn.line.width = Pt(width)
    if arrow:
        conn.line.end_arrowhead = True
    return conn


def panel(slide, x, y, w, h, title):
    box(slide, x, y, w, h, fill="white", edge="ink", lw=1.2, radius=True)
    text(slide, x + 0.05, y + 0.04, w - 0.10, 0.18, title, size=8.0,
         bold=True, align=PP_ALIGN.CENTER)


def node(slide, x, y, r=0.055):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - r), Inches(y - r),
                                 Inches(2 * r), Inches(2 * r))
    shp.fill.solid()
    shp.fill.fore_color.rgb = COL["white"]
    shp.line.color.rgb = COL["blue_edge"]
    shp.line.width = Pt(0.9)
    return shp


def graph_icon(slide, x, y, s=1.0):
    pts = [(x, y + 0.20 * s), (x + 0.26 * s, y), (x + 0.52 * s, y + 0.20 * s),
           (x + 0.16 * s, y + 0.48 * s), (x + 0.46 * s, y + 0.52 * s)]
    for a, b in [(0, 1), (1, 2), (0, 3), (3, 2), (3, 4), (2, 4)]:
        line(slide, pts[a][0], pts[a][1], pts[b][0], pts[b][1],
             color="blue_edge", width=0.7, arrow=False)
    for px, py in pts:
        node(slide, px, py, r=0.04 * s)


def mini_array(slide, x, y, rows=2, cols=3):
    box(slide, x, y, 0.42, 0.30, fill="white", edge="ink", radius=False, lw=0.7)
    for r in range(rows):
        for c in range(cols):
            shp = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x + 0.04 + c * 0.12),
                Inches(y + 0.04 + r * 0.11),
                Inches(0.08),
                Inches(0.07),
            )
            shp.fill.solid()
            shp.fill.fore_color.rgb = COL["white"]
            shp.line.color.rgb = COL["ink"]
            shp.line.width = Pt(0.5)


def draw():
    prs = Presentation()
    # Keep the original wide paper-figure size of Introduction.pptx.
    prs.slide_width = 6480175
    prs.slide_height = 2626995
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Three panels.
    panel(slide, 0.05, 0.07, 2.20, 2.65, "(a) TAG LLM Encoder Pipeline")
    panel(slide, 2.36, 0.07, 2.34, 2.65, "(b) Inefficiencies in Existing Execution")
    panel(slide, 4.82, 0.07, 2.20, 2.65, "(c) Proposed Framework: GRACE")
    line(slide, 2.25, 1.42, 2.36, 1.42, color="purple_edge", width=2.0)
    line(slide, 4.70, 1.42, 4.82, 1.42, color="purple_edge", width=2.0)

    # Panel A: pipeline and bottleneck.
    graph_icon(slide, 0.25, 0.47, s=1.0)
    text(slide, 0.17, 1.06, 0.72, 0.16, "TAG nodes", size=6.4, bold=True,
         align=PP_ALIGN.CENTER)
    box(slide, 0.92, 0.48, 0.68, 0.42, "LLM / TF\nEncoder", fill="blue",
        edge="blue_edge", size=6.3)
    box(slide, 1.70, 0.50, 0.42, 0.38, "GNN", fill="green", edge="green_edge",
        size=6.4)
    line(slide, 0.80, 0.72, 0.92, 0.72, width=0.9)
    line(slide, 1.60, 0.72, 1.70, 0.72, width=0.9)
    line(slide, 2.12, 0.72, 2.22, 0.72, width=0.9)
    text(slide, 0.18, 1.38, 1.92, 0.22, "Node-wise LLM encoding dominates", size=7.0,
         bold=True, align=PP_ALIGN.CENTER)
    box(slide, 0.28, 1.70, 1.45, 0.24, "Frontend LLM encoder", fill="blue",
        edge="blue_edge", size=5.7)
    box(slide, 1.76, 1.70, 0.23, 0.24, "GNN", fill="green",
        edge="green_edge", size=5.0)
    text(slide, 0.22, 2.14, 1.82, 0.30,
         "Graph context is available, yet ordinary TF accelerators encode every node independently.",
         size=5.5, color="muted", align=PP_ALIGN.CENTER)

    # Panel B: three observations/challenges.
    box(slide, 2.50, 0.43, 2.05, 0.56, fill="blue", edge="blue_edge")
    text(slide, 2.57, 0.50, 1.80, 0.14, "1. Graph locality reveals reuse", size=6.3,
         bold=True)
    graph_icon(slide, 2.62, 0.68, s=0.55)
    text(slide, 3.13, 0.66, 1.22, 0.22, "neighbor nodes are\ncloser than random pairs", size=5.2)
    box(slide, 4.28, 0.69, 0.20, 0.15, "~34%", fill="white", edge="blue_edge",
        size=4.8)

    box(slide, 2.50, 1.12, 2.05, 0.58, fill="orange", edge="orange_edge")
    text(slide, 2.57, 1.19, 1.83, 0.14, "2. Remaining encoder has BFPA boundary",
         size=6.1, bold=True)
    box(slide, 2.65, 1.42, 0.36, 0.16, "A3 fails", fill="red", edge="red_edge",
        size=4.6)
    box(slide, 3.18, 1.39, 0.42, 0.20, "A4 base", fill="yellow",
        edge="yellow_edge", size=4.8)
    box(slide, 3.80, 1.36, 0.46, 0.24, "A6 lift", fill="green",
        edge="green_edge", size=4.8)
    line(slide, 3.01, 1.50, 3.18, 1.50, color="orange_edge", width=0.8)
    line(slide, 3.60, 1.50, 3.80, 1.50, color="orange_edge", width=0.8)

    box(slide, 2.50, 1.84, 2.05, 0.58, fill="purple", edge="purple_edge")
    text(slide, 2.57, 1.91, 1.80, 0.14, "3. Reuse filtering and GEMMs diverge",
         size=6.1, bold=True)
    box(slide, 2.66, 2.16, 0.56, 0.18, "reuse filter", fill="gray",
        edge="gray_edge", size=4.8)
    box(slide, 3.62, 2.13, 0.62, 0.24, "dense GEMMs", fill="blue",
        edge="blue_edge", size=4.8)
    line(slide, 3.22, 2.25, 3.62, 2.25, color="purple_edge", width=0.8)

    # Panel C: GRACE framework.
    box(slide, 4.98, 0.40, 0.77, 1.80, "Memory-side\nreuse frontend", fill="green",
        edge="green_edge", size=5.6)
    box(slide, 5.08, 0.78, 0.56, 0.20, "lookup", fill="white",
        edge="green_edge", size=4.8)
    box(slide, 5.08, 1.07, 0.56, 0.20, "reuse filter", fill="white",
        edge="green_edge", size=4.6)
    box(slide, 5.08, 1.36, 0.56, 0.20, "prior results", fill="white",
        edge="green_edge", size=4.4)
    box(slide, 5.08, 1.65, 0.56, 0.20, "compute set", fill="white",
        edge="green_edge", size=4.4)

    box(slide, 6.02, 0.40, 0.78, 1.80, "BFP NPU\nbackend", fill="blue",
        edge="blue_edge", size=5.8)
    box(slide, 6.13, 0.82, 0.55, 0.18, "BFPA4 base", fill="yellow",
        edge="yellow_edge", size=4.8)
    box(slide, 6.13, 1.12, 0.55, 0.18, "BFPA6 lift", fill="orange",
        edge="orange_edge", size=4.8)
    mini_array(slide, 6.20, 1.43)
    box(slide, 6.13, 1.85, 0.55, 0.18, "TF GEMMs", fill="white",
        edge="blue_edge", size=4.8)
    line(slide, 5.75, 1.33, 6.02, 1.33, width=1.2)
    text(slide, 5.70, 1.43, 0.36, 0.20, "nodes+\nrisk tags", size=3.9,
         color="muted", align=PP_ALIGN.CENTER)

    box(slide, 4.98, 2.34, 1.82, 0.22,
        "1 encoder reuse   2 progressive BFP   3 heterogeneous execution",
        fill="purple", edge="purple_edge", size=4.7)

    prs.save(OUT)


if __name__ == "__main__":
    draw()
